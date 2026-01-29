#!/usr/bin/env python3
"""Create a single-slide summary template for revenue data."""

from pptx import Presentation
from pptx.util import Pt
import os

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUTPUT = os.path.join(ROOT, "templates", "Revenue_Summary_Template.potx")

prs = Presentation()
layout = prs.slide_layouts[1]  # Title and content
slide = prs.slides.add_slide(layout)

# Title
if slide.shapes.title:
    slide.shapes.title.text = "2025 Practice Revenue Summary"

# Content
for shape in slide.placeholders:
    if shape.placeholder_format.type == 7:  # BODY
        if shape.has_text_frame:
            tf = shape.text_frame
            tf.clear()
            
            fields = [
                ("Total Revenue", "{Total Revenue}"),
                ("Total Clients", "{Total Clients}"),
                ("Total Projects", "{Total Projects}"),
                ("Won Projects", "{Won Projects}"),
                ("", ""),
                ("Average Revenue per Project", "{Average Revenue per Project}"),
                ("Total GP $", "{Total GP $}"),
                ("Average GP %", "{Average GP %}"),
                ("Top Client", "{Top Client}"),
                ("Top Client Revenue", "{Top Client Revenue}"),
            ]
            
            for i, (label, placeholder) in enumerate(fields):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                if label:
                    p.text = f"{label}: {placeholder}"
                else:
                    p.text = ""
                p.font.size = Pt(16)
            break

prs.save(OUTPUT)
print(f"Created summary template: {OUTPUT}")
