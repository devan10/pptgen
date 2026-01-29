#!/usr/bin/env python3
"""Generate a template with placeholders for Practice Revenue data."""

from pptx import Presentation
from pptx.util import Inches, Pt
import os

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUTPUT = os.path.join(ROOT, "templates", "Revenue_Detail_Template.potx")

# Create presentation
prs = Presentation()

# Use Title and Content layout (usually layout 1)
layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(layout)

# Set title
if slide.shapes.title:
    slide.shapes.title.text = "Project: {Project Name}"

# Find and set content area
for shape in slide.placeholders:
    if shape.placeholder_format.type == 7:  # BODY placeholder
        if shape.has_text_frame:
            tf = shape.text_frame
            tf.clear()
            
            fields = [
                ("Client", "{Client}"),
                ("AE", "{AE}"),
                ("Practice Area", "{Practice Area}"),
                ("Status", "{Status}"),
                ("Date Signed", "{Date Signed}"),
                ("Resources", "{Resource(s)}"),
                ("Revenue", "${Rev}"),
                ("GP %", "{GP %}"),
                ("Hours", "{Hours}"),
                ("Comment", "{Comment}"),
            ]
            
            for i, (label, placeholder) in enumerate(fields):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"{label}: {placeholder}"
                p.font.size = Pt(14)
                p.level = 0
            break

prs.save(OUTPUT)
print(f"Created template: {OUTPUT}")

