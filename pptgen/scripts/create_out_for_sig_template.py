#!/usr/bin/env python3
"""Create an 'Out for Signature' projects table in the Insight template."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
import tempfile
import shutil
import zipfile

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
TEMPLATE_PATH = os.path.join(ROOT, "templates", "2026_Insight_PPT_Template.potx")
OUTPUT = os.path.join(ROOT, "templates", "2026_Insight_PPT_Template_OutForSig.potx")

# Convert .potx to .pptx temporarily
temp_dir = tempfile.mkdtemp()
temp_pptx = os.path.join(temp_dir, "temp.pptx")

with zipfile.ZipFile(TEMPLATE_PATH, 'r') as zip_in:
    with zipfile.ZipFile(temp_pptx, 'w') as zip_out:
        for item in zip_in.infolist():
            data = zip_in.read(item.filename)
            if item.filename == '[Content_Types].xml':
                data = data.replace(
                    b'application/vnd.openxmlformats-officedocument.presentationml.template.main+xml',
                    b'application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml'
                )
            zip_out.writestr(item, data)

# Load the insight template
prs = Presentation(temp_pptx)

# Remove existing slides
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[0]

# Use Content_Small 5 layout (to keep Insight template styling)
layout = prs.slide_layouts[14]  # Content_Small 5
slide = prs.slides.add_slide(layout)

# Set title
if slide.shapes.title:
    slide.shapes.title.text = "Projects Out for Signature"

# Remove all placeholders except title
for shape in list(slide.shapes):
    if shape != slide.shapes.title and shape.is_placeholder:
        sp = shape.element
        sp.getparent().remove(sp)

# Add table - adjusted for actual available space
rows, cols = 8, 4  # Header + 7 projects, 4 columns
left = Inches(0.5)
top = Inches(1.3)
width = Inches(9)
height = Inches(3.5)

table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
table = table_shape.table

# Set column widths
table.columns[0].width = Inches(2.5)  # Client
table.columns[1].width = Inches(3)    # Project Name
table.columns[2].width = Inches(1.5)  # Days to Sign
table.columns[3].width = Inches(1.5)  # Revenue

# Set row heights
for row_idx in range(rows):
    table.rows[row_idx].height = Inches(0.42)

# Add headers
headers = ['Client', 'Project Name', 'Days to Sign', 'Revenue']
for col_idx, header in enumerate(headers):
    cell = table.cell(0, col_idx)
    cell.text = header
    # Format header
    paragraph = cell.text_frame.paragraphs[0]
    paragraph.font.bold = True
    paragraph.font.size = Pt(10)
    paragraph.font.color.rgb = RGBColor(255, 255, 255)
    # Header background color (dark blue)
    fill = cell.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(31, 78, 121)

# Add placeholder rows with field names
for row_idx in range(1, rows):
    table.cell(row_idx, 0).text = f"{{Client_{row_idx}}}"
    table.cell(row_idx, 1).text = f"{{Project_Name_{row_idx}}}"
    table.cell(row_idx, 2).text = f"{{Days_to_Sign_{row_idx}}}"
    table.cell(row_idx, 3).text = f"{{Revenue_{row_idx}}}"
    
    # Format cells
    for col_idx in range(cols):
        cell = table.cell(row_idx, col_idx)
        # Enable word wrap for long text
        cell.text_frame.word_wrap = True
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.size = Pt(9)

prs.save(OUTPUT)
print(f"Created: {OUTPUT}")

# Cleanup
shutil.rmtree(temp_dir, ignore_errors=True)
