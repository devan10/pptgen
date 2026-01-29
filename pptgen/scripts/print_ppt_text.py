import sys
from pptx import Presentation
import os

def get_text_from_shape(shape):
    texts = []
    if hasattr(shape, 'text') and shape.text and shape.text.strip():
        texts.append(shape.text)
    elif hasattr(shape, 'text_frame') and shape.text_frame is not None:
        for para in shape.text_frame.paragraphs:
            if para.text and para.text.strip():
                texts.append(para.text)
    return texts


def main(path=None):
    if path is None:
        path = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'out', 'test_output.pptx')
    if not os.path.exists(path):
        print(f"File not found: {path}")
        sys.exit(2)
    prs = Presentation(path)
    for i, slide in enumerate(prs.slides, start=1):
        print(f"--- Slide {i} ---")
        any_text = False
        for shape in slide.shapes:
            texts = get_text_from_shape(shape)
            for t in texts:
                any_text = True
                print(t)
        if not any_text:
            print("(no text found)")

if __name__ == '__main__':
    arg = sys.argv[1] if len(sys.argv) > 1 else None
    main(arg)
