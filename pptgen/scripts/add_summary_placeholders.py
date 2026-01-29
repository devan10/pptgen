#!/usr/bin/env python3
"""Add summary placeholders to the Insight template."""

from pptx import Presentation
from pptx.util import Inches, Pt
import os
import tempfile
import shutil
import zipfile

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
TEMPLATE_PATH = os.path.join(ROOT, "templates", "2026_Insight_PPT_Template.potx")
OUTPUT = os.path.join(ROOT, "templates", "2026_Insight_PPT_Template_Summary.potx")

# Convert .potx to .pptx temporarily
temp_dir = tempfile.mkdtemp()
temp_pptx = os.path.join(temp_dir, "temp.pptx")

with zipfile.ZipFile(TEMPLATE_PATH, 'r') as zip_in:
    with zipfile.ZipFile(temp_pptx, 'w') as zip_out:
        for item in zip_in.infolist():
            data = zip_in.read(item.filename)
            if item.filename == '[Content_Types].xml':
                data = data.replace(
                    b'application/vnd.openxmlformats-officedocument.presentationml.template.main+xml',
                    b'application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml'
                )
            zip_out.writestr(item, data)

# Load the insight template
prs = Presentation(temp_pptx)

# Remove existing slides and use Content 5 layout
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[0]

# Use Content 5 layout (try multiple possible indices)
# Content 5 might be at index 4 or 5 depending on template
layout = None
for layout_idx in [4, 5, 6, 1]:
    if len(prs.slide_layouts) > layout_idx:
        test_layout = prs.slide_layouts[layout_idx]
        # Check if it has a body placeholder
        has_body = False
        for ph in test_layout.placeholders:
            if ph.placeholder_format.type == 7:  # BODY
                has_body = True
                break
        if has_body:
            layout = test_layout
            print(f"Using layout index {layout_idx}")
            break

if not layout:
    layout = prs.slide_layouts[1]

slide = prs.slides.add_slide(layout)

# Set title
if slide.shapes.title:
    slide.shapes.title.text = "2025 Practice Revenue Summary"

# Add content to the body placeholder
for shape in slide.placeholders:
    if shape.placeholder_format.type == 7:  # BODY placeholder
        if shape.has_text_frame:
            tf = shape.text_frame
            tf.clear()
            
            fields = [
                ("Total Revenue", "{Total Revenue}"),
                ("Total Projects", "{Total Projects}"),
                ("Total Clients", "{Total Clients}"),
                ("Won Projects", "{Won Projects}"),
                ("", ""),
                ("Average Revenue per Project", "{Average Revenue per Project}"),
                ("Total GP $", "{Total GP $}"),
                ("Average GP %", "{Average GP %}"),
                ("Top Client", "{Top Client}"),
                ("Top Client Revenue", "{Top Client Revenue}"),
            ]
            
            for i, (label, placeholder) in enumerate(fields):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                if label:
                    p.text = f"{label}: {placeholder}"
                else:
                    p.text = ""
                p.font.size = Pt(18)
            break

prs.save(OUTPUT)
print(f"Created: {OUTPUT}")

# Cleanup
shutil.rmtree(temp_dir, ignore_errors=True)

