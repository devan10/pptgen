from pptx import Presentation
from pptx.util import Inches, Pt
import subprocess
import os

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
DATA = os.path.join(ROOT, "data", "sample_data.csv")
TEMPLATE = os.path.join(ROOT, "scripts", "template.pptx")
OUTPUT_DIR = os.path.join(ROOT, "out")
OUTPUT = os.path.join(OUTPUT_DIR, "test_output.pptx")
GENERATOR = os.path.join(ROOT, "src", "generate_ppt.py")

# Create template with placeholders
prs = Presentation()
slide_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
# Add title placeholder text
if slide.shapes.title:
    slide.shapes.title.text = "{title}"
# Add a content textbox with placeholder
left = Inches(1)
top = Inches(1.5)
width = Inches(8)
height = Inches(4)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "{content}"
prs.save(TEMPLATE)
print(f"Wrote template: {TEMPLATE}")

# Ensure output dir exists
os.makedirs(OUTPUT_DIR, exist_ok=True)

# Run the generator script
cmd = ["python", GENERATOR, "--data", DATA, "--template", TEMPLATE, "--output", OUTPUT]
print("Running:", " ".join(cmd))
proc = subprocess.run(cmd, capture_output=True, text=True)
print(proc.stdout)
print(proc.stderr)

if os.path.exists(OUTPUT):
    print("Test succeeded, output written:", OUTPUT)
else:
    print("Test failed, output not found")
