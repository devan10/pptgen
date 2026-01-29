#!/usr/bin/env python3
"""Simple PPT generator.

Reads a CSV or Excel file and for each row adds a slide to the template.
Placeholders in template slides should use the form `{column_name}`.
"""

import argparse
import os
import shutil
import tempfile
import zipfile
import copy
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE


def replace_placeholders(slide, row):
    for shape in slide.shapes:
        # Handle tables
        if shape.has_table:
            table = shape.table
            for row_idx, table_row in enumerate(table.rows):
                for col_idx, cell in enumerate(table_row.cells):
                    text = cell.text
                    for col in row.index:
                        placeholder = "{" + str(col) + "}"
                        if placeholder in text:
                            val = "" if pd.isna(row[col]) else str(row[col])
                            text = text.replace(placeholder, val)
                    if text != cell.text:
                        cell.text_frame.clear()
                        cell.text_frame.text = text
        # Handle regular shapes
        elif not hasattr(shape, "text"):
            continue
        else:
            text = shape.text
            for col in row.index:
                placeholder = "{" + str(col) + "}"
                if placeholder in text:
                    val = "" if pd.isna(row[col]) else str(row[col])
                    text = text.replace(placeholder, val)
            if text != shape.text:
                try:
                    shape.text = text
                except Exception:
                    if hasattr(shape, "text_frame"):
                        shape.text_frame.clear()
                        shape.text_frame.text = text


def build_ppt(data_path, template_path, output_path, sheet=None):
    if data_path.lower().endswith(".csv"):
        df = pd.read_csv(data_path)
    else:
        df = pd.read_excel(data_path, sheet_name=sheet)

    if not os.path.exists(template_path):
        raise FileNotFoundError(f"Template not found: {template_path}")

    # Handle .potx (template) files by converting to .pptx for loading
    actual_template = template_path
    temp_dir = None
    if template_path.lower().endswith(".potx"):
        temp_dir = tempfile.mkdtemp()
        actual_template = os.path.join(temp_dir, "temp_template.pptx")
        # Convert .potx to .pptx by adjusting the content type in [Content_Types].xml
        with zipfile.ZipFile(template_path, 'r') as zip_in:
            with zipfile.ZipFile(actual_template, 'w') as zip_out:
                for item in zip_in.infolist():
                    data = zip_in.read(item.filename)
                    # Replace the presentation template content type with presentation type
                    if item.filename == '[Content_Types].xml':
                        data = data.replace(
                            b'application/vnd.openxmlformats-officedocument.presentationml.template.main+xml',
                            b'application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml'
                        )
                    zip_out.writestr(item, data)
    
    try:
        prs = Presentation(actual_template)
        
        # Get the first slide as the template and copy it for each row
        if len(prs.slides) == 0:
            raise ValueError("Template has no slides")
        
        template_slide = prs.slides[0]
        blank_slide_layout = prs.slide_layouts[6]  # Blank layout for copying

        for idx, (_, row) in enumerate(df.iterrows()):
            if idx == 0:
                # Use the first template slide for the first row
                slide = template_slide
            else:
                # Copy the template slide's structure
                # This is a workaround: add blank slide and copy shapes
                slide = prs.slides.add_slide(blank_slide_layout)
                for shape in template_slide.shapes:
                    if shape.has_text_frame:
                        el = shape.element
                        newel = copy.deepcopy(el)
                        slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
            
            replace_placeholders(slide, row)

        os.makedirs(os.path.dirname(output_path) or '.', exist_ok=True)
        prs.save(output_path)
        print(f"Wrote {output_path}")
    finally:
        if temp_dir:
            shutil.rmtree(temp_dir, ignore_errors=True)


def main():
    parser = argparse.ArgumentParser(description="Generate PPT from template and data file")
    parser.add_argument("--data", required=True, help="Path to Excel (.xlsx) or CSV file")
    parser.add_argument("--template", default="template.pptx", help="Path to template PPTX")
    parser.add_argument("--output", default="output.pptx", help="Output PPTX path")
    parser.add_argument("--sheet", default=None, help="Excel sheet name or index (optional)")
    args = parser.parse_args()

    build_ppt(args.data, args.template, args.output, sheet=args.sheet)


if __name__ == "__main__":
    main()
